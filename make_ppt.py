"""
SANJAYA EDGE — 7-Slide Hackathon Presentation Generator
Produces landscape PPTX (16:9) ready for PDF export.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, textwrap

BASE      = os.path.dirname(os.path.abspath(__file__))
SHOTS_DIR = os.path.join(BASE, "screenshots")
OUT_PPTX  = os.path.join(BASE, "SANJAYA_EDGE_Presentation.pptx")

# ─── Colour palette ───────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x08, 0x0D, 0x18)   # #080D18 deep navy
BG_CARD    = RGBColor(0x0E, 0x16, 0x28)   # #0E1628
BLUE       = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6 blue accent
BLUE_LIGHT = RGBColor(0x93, 0xC5, 0xFD)   # #93C5FD lighter blue
RED        = RGBColor(0xEF, 0x44, 0x44)   # #EF4444
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B
GREEN      = RGBColor(0x22, 0xC5, 0x5E)   # #22C55E
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT      = RGBColor(0xEE, 0xF2, 0xFF)   # #EEF2FF
MUTED      = RGBColor(0x8B, 0xA3, 0xC4)   # #8BA3C4
BORDER     = RGBColor(0x1E, 0x2F, 0x4A)   # #1E2F4A

# Slide dimensions — 16:9
W = Inches(13.33)
H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid() if fill else shape.fill.background()
    if fill:
        shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    if width and height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(width))
    elif height:
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 height=Inches(height))
    else:
        slide.shapes.add_picture(path, Inches(left), Inches(top))


def header_bar(slide, label="", show_logo=True):
    """Thin top bar with blue accent line."""
    add_rect(slide, 0, 0, 13.33, 0.55, fill=BG_CARD)
    add_rect(slide, 0, 0.52, 13.33, 0.04, fill=BLUE)
    if show_logo:
        add_text(slide, "SANJAYA EDGE", 0.3, 0.09, 3, 0.4,
                 font_size=13, bold=True, color=BLUE_LIGHT)
    if label:
        add_text(slide, label, 10.0, 0.09, 3.0, 0.4,
                 font_size=10, bold=False, color=MUTED,
                 align=PP_ALIGN.RIGHT)


def footer_bar(slide, left_txt="IIT Madras CoERS Road Safety Hackathon 2026 · DriveLegal Track",
               right_txt="github.com/Rajbharti06/SANJAYA-EDGE"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=BG_CARD)
    add_rect(slide, 0, 7.08, 13.33, 0.03, fill=BORDER)
    add_text(slide, left_txt,  0.3,  7.15, 8.0, 0.3,
             font_size=8, color=MUTED)
    add_text(slide, right_txt, 9.5,  7.15, 3.5, 0.3,
             font_size=8, color=BLUE_LIGHT, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, left, top, width, height,
                 dot_color=BLUE, text_color=LIGHT,
                 sub_color=MUTED, font_size=11, sub_size=9.5):
    """Render a list of (main, sub) bullet pairs."""
    y = top
    for (main, sub) in items:
        # Blue dot
        dot = slide.shapes.add_shape(1, Inches(left), Inches(y + 0.04),
                                     Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()

        add_text(slide, main, left + 0.17, y, width - 0.17, 0.25,
                 font_size=font_size, bold=True, color=text_color)
        if sub:
            add_text(slide, sub, left + 0.17, y + 0.23, width - 0.17, 0.22,
                     font_size=sub_size, color=sub_color)
        y += 0.5 if sub else 0.32


def badge(slide, text, left, top, width=1.4, height=0.30,
          fill=None, border=BLUE, text_color=BLUE_LIGHT, fs=9):
    fill_c = fill or BG_CARD
    add_rect(slide, left, top, width, height, fill=fill_c, line=border, line_w=Pt(1))
    add_text(slide, text, left + 0.05, top + 0.02, width - 0.1, height - 0.04,
             font_size=fs, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ─── SLIDE 1: WELCOME ─────────────────────────────────────────────────────────
def slide_welcome(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Background blue gradient band
    add_rect(sl, 0, 2.5, 13.33, 2.7, fill=RGBColor(0x0A, 0x14, 0x28))
    add_rect(sl, 0, 2.48, 13.33, 0.04, fill=BLUE)
    add_rect(sl, 0, 5.14, 13.33, 0.04, fill=BLUE)

    # Eye / target icon (simple shapes)
    cx, cy = 6.665, 1.4
    for r, c in [(0.55, BORDER), (0.38, RGBColor(0x1E,0x3A,0x8A)),
                 (0.22, BLUE), (0.08, BLUE_LIGHT)]:
        dot = sl.shapes.add_shape(9,  # oval
            Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()

    # Cross hairs
    for (x1,y1,x2,y2) in [(cx-0.7,cy,cx-0.28,cy),(cx+0.28,cy,cx+0.7,cy),
                           (cx,cy-0.7,cx,cy-0.28),(cx,cy+0.28,cx,cy+0.7)]:
        ln = sl.shapes.add_connector(1,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = BLUE
        ln.line.width = Pt(1.5)

    # Title
    add_text(sl, "SANJAYA  EDGE", 0, 2.65, 13.33, 1.0,
             font_size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(sl, "AI Road Safety Co-Pilot", 0, 3.55, 13.33, 0.55,
             font_size=22, bold=False, color=BLUE_LIGHT,
             align=PP_ALIGN.CENTER)

    # Tag line
    add_text(sl,
             "Real-Time Violation Detection  |  Legal Explainability  |  Citizen-Facing",
             0, 4.08, 13.33, 0.45,
             font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # Badges row
    badges = [("YOLOv8n", BLUE, BLUE_LIGHT),
              ("FastAPI", GREEN, RGBColor(0x86,0xEF,0xAC)),
              ("React + Vite", RGBColor(0x61,0xDA,0xFB), RGBColor(0xBA,0xF3,0xFF)),
              ("OpenCV HSV", AMBER, RGBColor(0xFC,0xD3,0x4D)),
              ("MVA 1988 + IPC 283", RGBColor(0xA3,0x5C,0xF7), RGBColor(0xC7,0xD2,0xFE))]
    bx = 1.8
    for (txt, bc, tc) in badges:
        badge(sl, txt, bx, 5.35, 1.8, 0.30, fill=BG_CARD, border=bc, text_color=tc, fs=9)
        bx += 1.95

    # Hackathon label
    add_text(sl, "IIT Madras CoERS Road Safety Hackathon 2026  |  DriveLegal Track",
             0, 5.9, 13.33, 0.4,
             font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(sl, "github.com/Rajbharti06/SANJAYA-EDGE",
             0, 6.3, 13.33, 0.35,
             font_size=10, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    return sl


# ─── SLIDE 2: PROBLEM STATEMENT ───────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    header_bar(sl, "02 / 07")
    footer_bar(sl)

    add_text(sl, "The Problem", 0.4, 0.65, 8, 0.6,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.20, 1.2, 0.05, fill=BLUE)

    # Stat cards row
    stats = [
        ("1,68,491",  "road crash deaths in India (2022)", RED),
        ("53%",       "of fatalities are 18-45 yr working age", AMBER),
        ("Top 3",     "causes: red-light jump, wrong-side, no helmet", RED),
        ("< 2%",      "violations ever result in a challan", AMBER),
    ]
    sx = 0.35
    for (val, lbl, col) in stats:
        add_rect(sl, sx, 1.4, 2.9, 1.5, fill=BG_CARD, line=BORDER)
        add_rect(sl, sx, 1.4, 2.9, 0.06, fill=col)
        add_text(sl, val, sx + 0.15, 1.55, 2.6, 0.65,
                 font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, sx + 0.1, 2.20, 2.7, 0.55,
                 font_size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
        sx += 3.15

    # Gap statement
    add_rect(sl, 0.35, 3.10, 12.6, 0.06, fill=BORDER)
    add_text(sl, "The Gap  —  Why Existing Systems Fail", 0.4, 3.25, 10, 0.4,
             font_size=16, bold=True, color=WHITE)

    gaps = [
        ("Reactive, not proactive",
         "Police enforce after the fact.  No real-time citizen warning."),
        ("Legally opaque",
         "Challans cite a section number — citizens don't know what law they broke or what fine applies."),
        ("No city-scale coverage",
         "Camera density is low.  Most violations go completely undetected."),
        ("No multi-violation awareness",
         "Existing tools detect one class.  Real intersections involve 4-5 simultaneous risks."),
    ]
    bullet_block(sl, gaps, 0.4, 3.75, 12.5, 3.0,
                 dot_color=RED, text_color=LIGHT, font_size=11.5)

    return sl


# ─── SLIDE 3: SOLUTION OVERVIEW ───────────────────────────────────────────────
def slide_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    header_bar(sl, "03 / 07")
    footer_bar(sl)

    add_text(sl, "Solution  —  SANJAYA EDGE", 0.4, 0.65, 10, 0.55,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.18, 1.2, 0.05, fill=BLUE)

    add_text(sl,
             "A real-time AI legal co-pilot that detects violations, cites the exact law, "
             "and announces the state-specific fine — proactively, before harm occurs.",
             0.4, 1.30, 12.5, 0.55,
             font_size=12, color=MUTED, italic=True)

    # 5 detector cards
    detectors = [
        ("D1", "Red Light Jump",
         "MVA Sec. 119",  "Rs. 1,000 / Rs. 2,000",  RED,
         "YOLOv8 + HSV 4-frame temporal"),
        ("D2", "Helmet Violation",
         "MVA Sec. 129",  "Rs. 1,000 + 3-mo. ban",   AMBER,
         "Custom head/helmet classifier"),
        ("D3", "Wrong-Side Driving",
         "MVA Sec. 184",  "Rs. 5,000",               RED,
         "Lucas-Kanade optical flow"),
        ("D4", "Traffic Blocking",
         "IPC Sec. 283",  "Rs. 5,000 + Rs. 2,000",   AMBER,
         "Vectorised IoU cluster >= 6"),
        ("D5", "Pedestrian in Lane",
         "MVA Sec. 283",  "Rs. 500",                  GREEN,
         "Aspect-ratio + zone filter"),
    ]

    dx = 0.3
    for (code, name, law, fine, col, method) in detectors:
        add_rect(sl, dx, 2.0, 2.45, 3.2, fill=BG_CARD, line=BORDER)
        add_rect(sl, dx, 2.0, 2.45, 0.07, fill=col)

        # Code badge
        add_rect(sl, dx + 0.12, 2.12, 0.55, 0.32, fill=col)
        add_text(sl, code, dx + 0.12, 2.12, 0.55, 0.32,
                 font_size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        add_text(sl, name, dx + 0.75, 2.12, 1.6, 0.32,
                 font_size=11, bold=True, color=LIGHT)
        add_text(sl, law,  dx + 0.12, 2.52, 2.2, 0.25,
                 font_size=9.5, color=BLUE_LIGHT)
        add_rect(sl, dx + 0.12, 2.82, 2.22, 0.28, fill=RGBColor(0x14,0x1E,0x38), line=col)
        add_text(sl, fine, dx + 0.17, 2.84, 2.12, 0.24,
                 font_size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, method, dx + 0.12, 3.20, 2.22, 0.4,
                 font_size=8.5, color=MUTED)

        dx += 2.57

    # What each alert delivers
    add_rect(sl, 0.3, 5.35, 12.7, 1.35, fill=BG_CARD, line=BORDER)
    add_rect(sl, 0.3, 5.35, 12.7, 0.05, fill=BLUE)
    add_text(sl, "Every Alert Delivers:", 0.55, 5.45, 5, 0.35,
             font_size=12, bold=True, color=WHITE)
    delivers = [
        "Violation name + exact MVA / IPC section",
        "National fine  +  TN / MH / KA state-specific amount",
        "Risk statement  (why it is dangerous)",
        "Voice announcement  (Web Speech API  |  en-IN)",
    ]
    dy = 5.88
    for item in delivers:
        dot = sl.shapes.add_shape(9, Inches(0.55), Inches(dy + 0.05),
                                  Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN
        dot.line.fill.background()
        add_text(sl, item, 0.70, dy, 11.5, 0.25,
                 font_size=10, color=LIGHT)
        dy += 0.28

    return sl


# ─── SLIDE 4: TECHNICAL ARCHITECTURE ─────────────────────────────────────────
def slide_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    header_bar(sl, "04 / 07")
    footer_bar(sl)

    add_text(sl, "Technical Architecture", 0.4, 0.65, 10, 0.55,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.18, 1.2, 0.05, fill=BLUE)

    # ── Backend box ──
    add_rect(sl, 0.3, 1.35, 5.0, 4.7, fill=BG_CARD, line=BORDER)
    add_rect(sl, 0.3, 1.35, 5.0, 0.06, fill=BLUE)
    add_text(sl, "BACKEND  (Python · FastAPI)", 0.45, 1.45, 4.7, 0.35,
             font_size=12, bold=True, color=BLUE_LIGHT)

    be_items = [
        ("YOLOv8n",         "COCO 80-class object detection"),
        ("OpenCV HSV",      "Traffic light colour analysis"),
        ("Lucas-Kanade",    "Optical flow wrong-side detection"),
        ("IoU Cluster",     "Traffic blocking gridlock detector"),
        ("Thread Executor", "Non-blocking async inference"),
        ("rules.json",      "MVA 1988 + IPC 283 legal database"),
    ]
    by = 1.9
    for (title, sub) in be_items:
        add_rect(sl, 0.45, by, 4.7, 0.52, fill=RGBColor(0x0A,0x12,0x22), line=BORDER)
        add_text(sl, title, 0.60, by + 0.04, 1.9, 0.24,
                 font_size=10, bold=True, color=LIGHT)
        add_text(sl, sub, 0.60, by + 0.27, 4.3, 0.22,
                 font_size=8.5, color=MUTED)
        by += 0.60

    # ── Arrow ──
    for yi in [2.8, 3.3, 3.8]:
        ln = sl.shapes.add_connector(1,
            Inches(5.3), Inches(yi), Inches(6.1), Inches(yi))
        ln.line.color.rgb = BLUE
        ln.line.width = Pt(2)
    add_text(sl, "WebSocket\nbase64 JPEG\n+ JSON", 5.28, 2.95, 0.9, 0.8,
             font_size=7.5, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # ── Frontend box ──
    add_rect(sl, 6.1, 1.35, 4.7, 4.7, fill=BG_CARD, line=BORDER)
    add_rect(sl, 6.1, 1.35, 4.7, 0.06, fill=RGBColor(0x61,0xDA,0xFB))
    add_text(sl, "FRONTEND  (React · Vite · Tailwind)", 6.25, 1.45, 4.5, 0.35,
             font_size=12, bold=True, color=RGBColor(0xBA,0xF3,0xFF))

    fe_items = [
        ("Live Feed",         "Annotated JPEG frames from WebSocket"),
        ("Event Log",         "Color-coded cards: Critical / High / Safe"),
        ("Stats Strip",       "Events · Critical · Confidence · Fine total"),
        ("Voice Co-Pilot",    "Web Speech API · en-IN · reads every alert"),
        ("5 Scenarios",       "Red light, helmet, wrong-side, blocking"),
        ("HUD Overlay",       "Speed, GPS, detector status, REC indicator"),
    ]
    fy = 1.9
    for (title, sub) in fe_items:
        add_rect(sl, 6.25, fy, 4.7, 0.52,
                 fill=RGBColor(0x0A,0x12,0x22), line=BORDER)
        add_text(sl, title, 6.40, fy + 0.04, 1.9, 0.24,
                 font_size=10, bold=True, color=LIGHT)
        add_text(sl, sub, 6.40, fy + 0.27, 4.3, 0.22,
                 font_size=8.5, color=MUTED)
        fy += 0.60

    # ── Performance table ──
    add_rect(sl, 11.0, 1.35, 2.1, 4.7, fill=BG_CARD, line=BORDER)
    add_rect(sl, 11.0, 1.35, 2.1, 0.06, fill=GREEN)
    add_text(sl, "Performance", 11.08, 1.45, 2.0, 0.35,
             font_size=10, bold=True, color=RGBColor(0x86,0xEF,0xAC))
    perf = [
        ("imgsz", "320px"),
        ("Skip", "3 frames"),
        ("Resize", "416px"),
        ("JPEG", "Q60 / 640px"),
        ("Pool", "2 threads"),
        ("Warmup", "startup"),
        ("CUDA", "auto"),
    ]
    py = 1.90
    for (k, v) in perf:
        add_text(sl, k, 11.08, py, 0.95, 0.24,
                 font_size=8.5, color=MUTED)
        add_text(sl, v, 12.0, py, 1.0, 0.24,
                 font_size=8.5, bold=True, color=BLUE_LIGHT,
                 align=PP_ALIGN.RIGHT)
        py += 0.54

    return sl


# ─── SLIDE 5: DEMO / SCREENSHOTS ──────────────────────────────────────────────
def slide_demo(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    header_bar(sl, "05 / 07")
    footer_bar(sl)

    add_text(sl, "Live Demo  —  Real Indian Traffic", 0.4, 0.65, 10, 0.55,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.18, 1.2, 0.05, fill=BLUE)

    # Layout: 2 rows x 3 cols of screenshots
    shots = [
        ("screenshots/01_01_Hello_everyone_-_this_is_SANJAYA_EDGE.png",
         "YOLOv8n live scanning", BLUE),
        ("screenshots/01_03_DETECTOR_1_Red_Light_Jump_MVA_Section_11.png",
         "Red Light Detection", RED),
        ("screenshots/02_03_Bare_head_detected_on_motorcyclist.png",
         "Helmet Violation", AMBER),
        ("screenshots/03_03_Vehicle_moving_140_degrees_against_domin.png",
         "Wrong-Side Driving", RED),
        ("screenshots/04_02_Dense_vehicle_cluster_detected_6_vehicle.png",
         "Traffic Blocking", AMBER),
        ("screenshots/04_03_ARCHITECTURE_FastAPI_React_WebSocket_str.png",
         "Architecture Overview", BLUE),
    ]

    cols = 3
    img_w = 3.9
    img_h = 2.2
    gap_x = 0.35
    gap_y = 0.25
    start_x = 0.3
    start_y = 1.35

    for idx, (path, caption, col) in enumerate(shots):
        row = idx // cols
        col_i = idx % cols
        x = start_x + col_i * (img_w + gap_x)
        y = start_y + row * (img_h + gap_y + 0.32)

        full = os.path.join(BASE, path)
        add_rect(sl, x, y, img_w, img_h, fill=BG_CARD, line=BORDER)
        add_image(sl, full, x + 0.04, y + 0.04, width=img_w - 0.08,
                  height=img_h - 0.08)

        # Caption bar under image
        add_rect(sl, x, y + img_h, img_w, 0.28, fill=BG_CARD, line=BORDER)
        dot = sl.shapes.add_shape(9, Inches(x + 0.08), Inches(y + img_h + 0.09),
                                  Inches(0.09), Inches(0.09))
        dot.fill.solid(); dot.fill.fore_color.rgb = col
        dot.line.fill.background()
        add_text(sl, caption, x + 0.22, y + img_h + 0.02, img_w - 0.28, 0.24,
                 font_size=9, bold=True, color=LIGHT)

    return sl


# ─── SLIDE 6: INNOVATION & IMPACT ─────────────────────────────────────────────
def slide_impact(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    header_bar(sl, "06 / 07")
    footer_bar(sl)

    add_text(sl, "Innovation  &  Impact", 0.4, 0.65, 10, 0.55,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.4, 1.18, 1.2, 0.05, fill=BLUE)

    # Left column — innovation
    add_rect(sl, 0.3, 1.35, 6.1, 5.35, fill=BG_CARD, line=BORDER)
    add_rect(sl, 0.3, 1.35, 6.1, 0.06, fill=BLUE)
    add_text(sl, "What Makes It Different", 0.48, 1.45, 5.8, 0.35,
             font_size=13, bold=True, color=BLUE_LIGHT)

    innovations = [
        ("Temporal Consensus",
         "4-frame HSV filter eliminates false red-light positives from sunlight glare."),
        ("Optical Flow Detection",
         "Lucas-Kanade flow on vehicle centres flags wrong-side with no GPS needed."),
        ("Citizen-Facing Legal DB",
         "Every alert names the exact law section + fine per state (TN/MH/KA)."),
        ("Thread-Pool Inference",
         "YOLO runs in executor — async event loop never blocks; WebSocket stays smooth."),
        ("Voice Co-Pilot",
         "Web Speech API (en-IN) announces each violation aloud within 200ms."),
        ("5 Detectors Simultaneously",
         "Red light + helmet + wrong-side + blocking + pedestrian all run in one pass."),
    ]
    bullet_block(sl, innovations, 0.48, 1.90, 5.8, 4.6,
                 dot_color=BLUE, text_color=LIGHT, font_size=10.5)

    # Right column — impact
    add_rect(sl, 6.6, 1.35, 6.4, 5.35, fill=BG_CARD, line=BORDER)
    add_rect(sl, 6.6, 1.35, 6.4, 0.06, fill=GREEN)
    add_text(sl, "Potential Impact", 6.78, 1.45, 6.1, 0.35,
             font_size=13, bold=True, color=RGBColor(0x86,0xEF,0xAC))

    impact_items = [
        ("Scalable to any camera",
         "Dashcam, CCTV, traffic pole cam — no hardware lock-in."),
        ("Zero police dependency",
         "Citizens receive legal education at point-of-violation, not after."),
        ("State-specific fines",
         "Covers TN, MH, KA.  Database extensible to all 28 states."),
        ("Proactive deterrence",
         "Real-time voice warning before the next violation, not after the first."),
        ("Open source",
         "MIT-licensed.  Deployable by any smart-city or transport authority."),
        ("Roadmap: seatbelt, phone use",
         "3 more detectors ready to add.  Legal DB already has entries."),
    ]
    bullet_block(sl, impact_items, 6.78, 1.90, 6.1, 4.6,
                 dot_color=GREEN, text_color=LIGHT, font_size=10.5)

    return sl


# ─── SLIDE 7: THANK YOU ───────────────────────────────────────────────────────
def slide_thankyou(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Background band
    add_rect(sl, 0, 2.4, 13.33, 2.9, fill=RGBColor(0x0A, 0x14, 0x28))
    add_rect(sl, 0, 2.38, 13.33, 0.04, fill=BLUE)
    add_rect(sl, 0, 5.24, 13.33, 0.04, fill=BLUE)

    add_text(sl, "Thank You", 0, 2.55, 13.33, 0.85,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(sl, "SANJAYA EDGE  —  Making Every Road Safer, One Detection at a Time",
             0, 3.38, 13.33, 0.50,
             font_size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

    # Info cards
    info = [
        ("Code Repository",  "github.com/Rajbharti06/SANJAYA-EDGE", BLUE),
        ("Tech Stack",       "YOLOv8n · FastAPI · React · OpenCV · Python 3.11", MUTED),
        ("Legal Coverage",   "MVA 1988 · IPC 283 · TN / MH / KA State Rules", AMBER),
        ("Hackathon",        "IIT Madras CoERS Road Safety Hackathon 2026 · DriveLegal", GREEN),
    ]
    ix = 0.55
    for (label, val, col) in info:
        add_rect(sl, ix, 5.45, 2.9, 0.85, fill=BG_CARD, line=BORDER)
        add_rect(sl, ix, 5.45, 2.9, 0.05, fill=col)
        add_text(sl, label, ix + 0.12, 5.55, 2.7, 0.28,
                 font_size=9, color=MUTED, bold=False)
        add_text(sl, val,   ix + 0.12, 5.80, 2.7, 0.42,
                 font_size=9.5, bold=True, color=LIGHT)
        ix += 3.10

    add_text(sl, "Submitted for IIT Madras CoERS Road Safety Hackathon 2026",
             0, 6.65, 13.33, 0.35,
             font_size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    return sl


# ─── BUILD ────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()

    print("[ppt] Slide 1 — Welcome")
    slide_welcome(prs)
    print("[ppt] Slide 2 — Problem Statement")
    slide_problem(prs)
    print("[ppt] Slide 3 — Solution Overview")
    slide_solution(prs)
    print("[ppt] Slide 4 — Technical Architecture")
    slide_architecture(prs)
    print("[ppt] Slide 5 — Demo / Screenshots")
    slide_demo(prs)
    print("[ppt] Slide 6 — Innovation & Impact")
    slide_impact(prs)
    print("[ppt] Slide 7 — Thank You")
    slide_thankyou(prs)

    prs.save(OUT_PPTX)
    size_mb = os.path.getsize(OUT_PPTX) / (1024 * 1024)
    print(f"\n[ppt] Saved: {OUT_PPTX}  ({size_mb:.1f} MB)")
    print("[ppt] Open in PowerPoint and export as PDF (File > Export > PDF) for final submission.")


if __name__ == "__main__":
    main()
